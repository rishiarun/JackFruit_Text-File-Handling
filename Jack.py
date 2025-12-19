"""
text_string_tools_customtk_all_formats.py

CustomTkinter GUI app with:
  1) Word Frequency (multi-format file browser):
       - Supports: .txt, .md, .log, .csv, .html, .htm, .docx, .pdf, .pptx
       - Ignores images and image-only content (scanned PDFs will appear empty)
       - Aligned "word.............count" output with monospace font
  2) Palindrome Checker
  3) Caesar Cipher (encrypt/decrypt)

Features:
  - Dark / Light / System appearance switch
"""

import os
import re
import string
from collections import Counter

import customtkinter as ctk
from tkinter import messagebox, filedialog

# External libraries for document formats
import docx           # python-docx
import PyPDF2         # PyPDF2
from pptx import Presentation  # python-pptx


# ---------------- LOGIC FUNCTIONS ---------------- #

def normalize_text_for_words(text: str) -> str:
    """Lowercase and replace punctuation with spaces."""
    trans_table = str.maketrans(string.punctuation, " " * len(string.punctuation))
    return text.translate(trans_table).lower()


def compute_word_frequency_from_text(text: str) -> Counter:
    """Return a Counter mapping words to frequencies from raw text."""
    normalized = normalize_text_for_words(text)
    words = normalized.split()
    return Counter(words)


def is_palindrome_core(s: str) -> bool:
    """Check if string is palindrome ignoring non-alnum and case."""
    cleaned_chars = []
    for ch in s:
        if ch.isalnum():
            cleaned_chars.append(ch.lower())
    cleaned = "".join(cleaned_chars)
    return cleaned == cleaned[::-1]


def caesar_transform(text: str, shift: int) -> str:
    """Apply Caesar shift to letters, preserving case and non-letters."""
    result_chars = []
    for ch in text:
        if ch.isalpha():
            if ch.isupper():
                base = ord('A')
            else:
                base = ord('a')
            new_ch = chr((ord(ch) - base + shift) % 26 + base)
            result_chars.append(new_ch)
        else:
            result_chars.append(ch)
    return "".join(result_chars)


# ---------------- FILE TEXT EXTRACTION ---------------- #

def extract_text_from_txt_like(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_text_from_html(path: str) -> str:
    # basic: strip HTML tags, including <img> etc.
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    # remove tags
    text = re.sub(r"<[^>]+>", " ", html)
    return text


def extract_text_from_docx(path: str) -> str:
    d = docx.Document(path)
    parts = []
    for para in d.paragraphs:
        if para.text:
            parts.append(para.text)
    return "\n".join(parts)


def extract_text_from_pdf(path: str) -> str:
    """Extract text from a PDF using PyPDF2; images are ignored by default."""
    text_parts = []
    with open(path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n".join(text_parts)


def extract_text_from_pptx(path: str) -> str:
    """Extract text from PowerPoint; ignores images."""
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text:
                    parts.append(shape.text)
            elif hasattr(shape, "text_frame"):
                # older style but generally covered by shape.text
                if shape.text_frame is not None:
                    for para in shape.text_frame.paragraphs:
                        run_texts = []
                        for run in para.runs:
                            run_texts.append(run.text)
                        para_text = "".join(run_texts)
                        if para_text:
                            parts.append(para_text)
    return "\n".join(parts)


def extract_text_generic(path: str) -> str:
    """
    Extract text from various file types.

    Supported:
      - .txt, .md, .log, .csv
      - .html, .htm (tags stripped; images ignored)
      - .docx
      - .pdf
      - .pptx

    Image-only content (e.g., scanned PDFs) will result in empty text.
    """
    _, ext = os.path.splitext(path)
    ext = ext.lower()

    if ext in (".txt", ".md", ".log", ".csv"):
        return extract_text_from_txt_like(path)
    if ext in (".html", ".htm"):
        return extract_text_from_html(path)
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext == ".pptx":
        return extract_text_from_pptx(path)

    # Unsupported or binary/image-only
    raise ValueError(f"Unsupported file type: {ext}")


# ---------------- MAIN APP CLASS ---------------- #

class TextToolsApp(ctk.CTk):

    def __init__(self):
        super().__init__()

        # basic window setup
        self.title("Text & String Tools (customtkinter)")
        self.geometry("1000x650")

        # appearance + theme
        ctk.set_appearance_mode("Dark")      # "Dark", "Light", or "System"
        ctk.set_default_color_theme("blue")  # "blue", "dark-blue", "green"

        # monospace font (used for text areas)
        self.mono_font = ("Consolas", 11)

        # build UI
        self.create_top_bar()
        self.create_tabs()

    # ---------------- TOP BAR (TITLE + THEME SWITCH) ---------------- #

    def create_top_bar(self):
        top_frame = ctk.CTkFrame(self, corner_radius=0)
        top_frame.pack(side="top", fill="x")

        title_label = ctk.CTkLabel(
            top_frame,
            text="Text & String Tools",
            font=("Segoe UI", 18, "bold")
        )
        title_label.pack(side="left", padx=20, pady=10)

        # appearance mode selector (Dark / Light / System)
        mode_label = ctk.CTkLabel(top_frame, text="Theme:")
        mode_label.pack(side="right", padx=(0, 5), pady=10)

        self.appearance_option = ctk.CTkOptionMenu(
            top_frame,
            values=["Dark", "Light", "System"],
            command=self.change_appearance_mode
        )
        self.appearance_option.set("Dark")
        self.appearance_option.pack(side="right", padx=10, pady=10)

    def change_appearance_mode(self, new_mode: str):
        ctk.set_appearance_mode(new_mode)

    # ---------------- TABS ---------------- #

    def create_tabs(self):
        # Tab view holds the three tools
        self.tabview = ctk.CTkTabview(self)
        self.tabview.pack(expand=True, fill="both", padx=10, pady=10)

        # create three tabs
        self.tab_word = self.tabview.add("Word Frequency")
        self.tab_palindrome = self.tabview.add("Palindrome Checker")
        self.tab_caesar = self.tabview.add("Caesar Cipher")

        self.build_word_tab()
        self.build_palindrome_tab()
        self.build_caesar_tab()

    # ---------------- WORD FREQUENCY TAB ---------------- #

    def build_word_tab(self):
        container = ctk.CTkFrame(self.tab_word)
        container.pack(expand=True, fill="both", padx=10, pady=10)

        # row: file path entry + browse button
        file_row = ctk.CTkFrame(container)
        file_row.pack(fill="x", pady=(0, 8))

        file_label = ctk.CTkLabel(file_row, text="Selected file:")
        file_label.pack(side="left", padx=(0, 5))

        self.wf_file_path_var = ctk.StringVar()
        self.wf_file_entry = ctk.CTkEntry(
            file_row,
            textvariable=self.wf_file_path_var,
            width=600
        )
        self.wf_file_entry.pack(side="left", fill="x", expand=True)

        browse_button = ctk.CTkButton(
            file_row,
            text="Browse...",
            command=self.wf_browse_file
        )
        browse_button.pack(side="left", padx=(10, 0))

        # row: Analyze + Clear buttons
        button_row = ctk.CTkFrame(container)
        button_row.pack(fill="x", pady=(0, 8))

        analyze_button = ctk.CTkButton(
            button_row,
            text="Analyze File",
            command=self.wf_analyze_file
        )
        analyze_button.pack(side="left")

        clear_button = ctk.CTkButton(
            button_row,
            text="Clear Output",
            fg_color="gray",
            hover_color="#666666",
            command=self.wf_clear_output
        )
        clear_button.pack(side="left", padx=(8, 0))

        # textbox for output
        self.wf_output = ctk.CTkTextbox(
            container,
            wrap="none"  # keep columns aligned horizontally
        )
        self.wf_output.configure(font=self.mono_font)
        self.wf_output.pack(expand=True, fill="both", pady=(5, 0))

    def wf_browse_file(self):
        file_path = filedialog.askopenfilename(
            title="Select a file",
            filetypes=[
                ("Supported files",
                 "*.txt *.md *.log *.csv *.html *.htm *.docx *.pdf *.pptx"),
                ("Text files", "*.txt *.md *.log *.csv"),
                ("Documents", "*.docx *.pdf *.pptx"),
                ("HTML files", "*.html *.htm"),
                ("All files", "*.*"),
            ]
        )
        if file_path:
            self.wf_file_path_var.set(file_path)

    def wf_analyze_file(self):
        path = self.wf_file_path_var.get().strip()
        if path == "":
            messagebox.showwarning("No file", "Please select a file using Browse.")
            return

        try:
            raw_text = extract_text_generic(path)
        except ValueError as ve:
            messagebox.showerror("Unsupported file", str(ve))
            return
        except Exception as e:
            messagebox.showerror("File error", f"Could not open or parse file:\n{e}")
            return

        if raw_text.strip() == "":
            messagebox.showwarning(
                "No text found",
                "No extractable text was found.\n"
                "If this is a scanned PDF or image-only document, it cannot be analyzed."
            )
            self.wf_output.delete("1.0", "end")
            self.wf_output.insert("end", "No text found in document.\n")
            return

        counter = compute_word_frequency_from_text(raw_text)
        if not counter:
            self.wf_output.delete("1.0", "end")
            self.wf_output.insert("end", "No words found after processing.\n")
            return

        # pretty aligned output with dots: word.....count
        max_word_len = 0
        for word in counter.keys():
            if len(word) > max_word_len:
                max_word_len = len(word)

        # where count column starts (in characters)
        count_column = max_word_len + 30  # tweak for more/less dots

        lines = []
        lines.append(f"Word frequency for file: {path}")
        lines.append("")

        for word, count in counter.most_common():
            dots_needed = count_column - len(word) - 1  # 1 space after word
            if dots_needed < 1:
                dots_needed = 1
            dots = "." * dots_needed
            line = f"{word} {dots} {count}"
            lines.append(line)

        self.wf_output.delete("1.0", "end")
        self.wf_output.insert("end", "\n".join(lines))

    def wf_clear_output(self):
        self.wf_output.delete("1.0", "end")
        self.wf_file_path_var.set("")

    # ---------------- PALINDROME TAB ---------------- #

    def build_palindrome_tab(self):
        container = ctk.CTkFrame(self.tab_palindrome)
        container.pack(expand=True, fill="both", padx=10, pady=10)

        label = ctk.CTkLabel(
            container,
            text="Enter text to check:",
            anchor="w"
        )
        label.pack(fill="x", pady=(0, 5))

        self.pal_input = ctk.CTkEntry(
            container,
            width=600
        )
        self.pal_input.pack(fill="x", pady=(0, 10))

        button_row = ctk.CTkFrame(container)
        button_row.pack(fill="x", pady=(0, 8))

        check_button = ctk.CTkButton(
            button_row,
            text="Check",
            command=self.pal_check
        )
        check_button.pack(side="left")

        clear_button = ctk.CTkButton(
            button_row,
            text="Clear",
            fg_color="gray",
            hover_color="#666666",
            command=self.pal_clear
        )
        clear_button.pack(side="left", padx=(8, 0))

        self.pal_result_label = ctk.CTkLabel(
            container,
            text="",
            font=("Segoe UI", 14, "bold")
        )
        self.pal_result_label.pack(anchor="w", pady=(10, 0))

    def pal_check(self):
        text = self.pal_input.get()
        if text.strip() == "":
            messagebox.showwarning("No input", "Please enter some text to check.")
            return

        if is_palindrome_core(text):
            self.pal_result_label.configure(
                text="Palindrome ✓",
                text_color="#39D353"  # green
            )
        else:
            self.pal_result_label.configure(
                text="Not a palindrome ✗",
                text_color="#FF6B6B"  # red
            )

    def pal_clear(self):
        self.pal_input.delete(0, "end")
        self.pal_result_label.configure(text="", text_color=None)

    # ---------------- CAESAR TAB ---------------- #

    def build_caesar_tab(self):
        container = ctk.CTkFrame(self.tab_caesar)
        container.pack(expand=True, fill="both", padx=10, pady=10)

        in_label = ctk.CTkLabel(
            container,
            text="Input text:",
            anchor="w"
        )
        in_label.pack(fill="x", pady=(0, 5))

        self.c_input_text = ctk.CTkTextbox(
            container,
            height=160
        )
        self.c_input_text.configure(font=self.mono_font)
        self.c_input_text.pack(fill="both", expand=False, pady=(0, 10))

        shift_row = ctk.CTkFrame(container)
        shift_row.pack(fill="x", pady=(0, 8))

        shift_label = ctk.CTkLabel(
            shift_row,
            text="Shift (integer):"
        )
        shift_label.pack(side="left", padx=(0, 5))

        self.c_shift_var = ctk.StringVar(value="3")
        self.c_shift_entry = ctk.CTkEntry(
            shift_row,
            textvariable=self.c_shift_var,
            width=60
        )
        self.c_shift_entry.pack(side="left")

        button_row = ctk.CTkFrame(container)
        button_row.pack(fill="x", pady=(0, 8))

        encrypt_button = ctk.CTkButton(
            button_row,
            text="Encrypt →",
            command=self.c_encrypt
        )
        encrypt_button.pack(side="left")

        decrypt_button = ctk.CTkButton(
            button_row,
            text="Decrypt ←",
            fg_color="#444444",
            hover_color="#666666",
            command=self.c_decrypt
        )
        decrypt_button.pack(side="left", padx=(8, 0))

        clear_button = ctk.CTkButton(
            button_row,
            text="Clear",
            fg_color="gray",
            hover_color="#666666",
            command=self.c_clear
        )
        clear_button.pack(side="left", padx=(8, 0))

        out_label = ctk.CTkLabel(
            container,
            text="Result:",
            anchor="w"
        )
        out_label.pack(fill="x", pady=(10, 5))

        self.c_output_text = ctk.CTkTextbox(
            container,
            height=160
        )
        self.c_output_text.configure(font=self.mono_font)
        self.c_output_text.pack(fill="both", expand=True)

    def get_shift_value(self):
        raw = self.c_shift_var.get().strip()
        try:
            value = int(raw)
            return value % 26
        except ValueError:
            messagebox.showerror("Invalid shift", "Shift must be an integer.")
            return None

    def c_encrypt(self):
        shift = self.get_shift_value()
        if shift is None:
            return

        text = self.c_input_text.get("1.0", "end").rstrip("\n")
        if text == "":
            messagebox.showwarning("No input", "Enter some text to encrypt.")
            return

        result = caesar_transform(text, shift)
        self.c_output_text.delete("1.0", "end")
        self.c_output_text.insert("end", result)

    def c_decrypt(self):
        shift = self.get_shift_value()
        if shift is None:
            return

        text = self.c_input_text.get("1.0", "end").rstrip("\n")
        if text == "":
            messagebox.showwarning("No input", "Enter some text to decrypt.")
            return

        result = caesar_transform(text, -shift)
        self.c_output_text.delete("1.0", "end")
        self.c_output_text.insert("end", result)

    def c_clear(self):
        self.c_input_text.delete("1.0", "end")
        self.c_output_text.delete("1.0", "end")
        self.c_shift_var.set("3")


# ---------------- RUN APP ---------------- #

if __name__ == "__main__":
    app = TextToolsApp()
    app.mainloop()
